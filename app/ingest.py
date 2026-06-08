"""
Document ingestion layer.

Takes raw uploaded bytes and extracts plain text, tables, and metadata from a
range of file formats. This layer is fully offline and deterministic — it does
not interpret the content, it only normalises it into an `ExtractedDoc` that the
analysis layer (app/analysis.py) can reason over.

Supported: PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx), CSV/TSV,
plain text / markdown / logs / JSON, and images (metadata only).
"""

from __future__ import annotations

import csv
import io
import json
from dataclasses import dataclass, field
from pathlib import Path

# Safety caps so a huge upload can't exhaust memory.
MAX_TABLE_ROWS = 5000
MAX_TEXT_CHARS = 2_000_000

TEXT_EXTS = {".txt", ".md", ".markdown", ".log", ".rst"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tif", ".tiff"}


@dataclass
class ExtractedDoc:
    """Normalised representation of one ingested file."""

    filename: str
    filetype: str  # pdf | docx | xlsx | pptx | csv | text | json | image | unknown
    text: str = ""
    tables: list[list[list[str]]] = field(default_factory=list)  # tables -> rows -> cells
    metadata: dict = field(default_factory=dict)
    error: str | None = None
    # For images only: base64 data + media type, so the analysis layer can send
    # them to the Claude vision API. Empty for non-image files.
    image_b64: str | None = None
    media_type: str | None = None

    @property
    def word_count(self) -> int:
        return len(self.text.split())


def extract(filename: str, data: bytes) -> ExtractedDoc:
    """Dispatch to the right extractor based on file extension."""
    ext = Path(filename).suffix.lower()
    try:
        if ext == ".pdf":
            return _extract_pdf(filename, data)
        if ext == ".docx":
            return _extract_docx(filename, data)
        if ext == ".xlsx":
            return _extract_xlsx(filename, data)
        if ext == ".pptx":
            return _extract_pptx(filename, data)
        if ext in (".csv", ".tsv"):
            return _extract_csv(filename, data, delimiter="\t" if ext == ".tsv" else ",")
        if ext == ".json":
            return _extract_json(filename, data)
        if ext in TEXT_EXTS:
            return _extract_text(filename, data, "text")
        if ext in IMAGE_EXTS:
            return _extract_image(filename, data)
        # Fallback: try to read it as UTF-8 text.
        return _extract_text(filename, data, "unknown")
    except Exception as exc:  # pragma: no cover - defensive: never crash the batch
        return ExtractedDoc(filename=filename, filetype=ext.lstrip(".") or "unknown", error=str(exc))


def _truncate(text: str) -> str:
    return text if len(text) <= MAX_TEXT_CHARS else text[:MAX_TEXT_CHARS] + "\n...[truncated]"


def _extract_pdf(filename: str, data: bytes) -> ExtractedDoc:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = [(p.extract_text() or "") for p in reader.pages]
    text = "\n".join(pages)
    return ExtractedDoc(
        filename=filename,
        filetype="pdf",
        text=_truncate(text),
        metadata={"pages": len(reader.pages)},
    )


def _extract_docx(filename: str, data: bytes) -> ExtractedDoc:
    import docx

    document = docx.Document(io.BytesIO(data))
    paras = [p.text for p in document.paragraphs if p.text.strip()]
    tables: list[list[list[str]]] = []
    for tbl in document.tables:
        rows = [[cell.text for cell in row.cells] for row in tbl.rows[:MAX_TABLE_ROWS]]
        if rows:
            tables.append(rows)
    return ExtractedDoc(
        filename=filename,
        filetype="docx",
        text=_truncate("\n".join(paras)),
        tables=tables,
        metadata={"paragraphs": len(paras), "tables": len(tables)},
    )


def _extract_xlsx(filename: str, data: bytes) -> ExtractedDoc:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    tables: list[list[list[str]]] = []
    lines: list[str] = []
    for ws in wb.worksheets:
        rows: list[list[str]] = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= MAX_TABLE_ROWS:
                break
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                rows.append(cells)
                lines.append(" | ".join(cells))
        if rows:
            tables.append(rows)
    wb.close()
    return ExtractedDoc(
        filename=filename,
        filetype="xlsx",
        text=_truncate("\n".join(lines)),
        tables=tables,
        metadata={"sheets": [ws.title for ws in wb.worksheets] if wb.worksheets else [],
                  "sheet_count": len(tables)},
    )


def _extract_pptx(filename: str, data: bytes) -> ExtractedDoc:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    chunks: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        slide_text.append(line)
        if slide_text:
            chunks.append(f"[Slide {idx}]\n" + "\n".join(slide_text))
    return ExtractedDoc(
        filename=filename,
        filetype="pptx",
        text=_truncate("\n\n".join(chunks)),
        metadata={"slides": len(prs.slides._sldIdLst)},
    )


def _extract_csv(filename: str, data: bytes, delimiter: str = ",") -> ExtractedDoc:
    text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    rows = [row for _, row in zip(range(MAX_TABLE_ROWS), reader)]
    lines = [" | ".join(r) for r in rows]
    return ExtractedDoc(
        filename=filename,
        filetype="csv",
        text=_truncate("\n".join(lines)),
        tables=[rows] if rows else [],
        metadata={"rows": len(rows), "columns": len(rows[0]) if rows else 0},
    )


def _extract_json(filename: str, data: bytes) -> ExtractedDoc:
    text = data.decode("utf-8", errors="replace")
    try:
        obj = json.loads(text)
        pretty = json.dumps(obj, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        pretty = text
    return ExtractedDoc(
        filename=filename, filetype="json", text=_truncate(pretty), metadata={"bytes": len(data)}
    )


def _extract_text(filename: str, data: bytes, filetype: str) -> ExtractedDoc:
    text = data.decode("utf-8", errors="replace")
    return ExtractedDoc(
        filename=filename, filetype=filetype, text=_truncate(text), metadata={"bytes": len(data)}
    )


_PIL_TO_MEDIA_TYPE = {
    "JPEG": "image/jpeg",
    "PNG": "image/png",
    "GIF": "image/gif",
    "WEBP": "image/webp",
}


def _extract_image(filename: str, data: bytes) -> ExtractedDoc:
    import base64

    from PIL import Image

    with Image.open(io.BytesIO(data)) as img:
        meta = {"format": img.format, "mode": img.mode, "width": img.width, "height": img.height}
        fmt = (img.format or "").upper()

    note = (
        f"[Image: {meta['width']}x{meta['height']} {meta.get('format')}. "
        "Visual content analysis requires the Claude vision API (offline mock mode reads metadata only).]"
    )
    # Carry the raw bytes (base64) so the analysis layer can pass it to Claude's
    # vision API. Claude supports JPEG/PNG/GIF/WEBP; skip others (e.g. TIFF/BMP).
    media_type = _PIL_TO_MEDIA_TYPE.get(fmt)
    return ExtractedDoc(
        filename=filename,
        filetype="image",
        text=note,
        metadata=meta,
        image_b64=base64.standard_b64encode(data).decode("ascii") if media_type else None,
        media_type=media_type,
    )
